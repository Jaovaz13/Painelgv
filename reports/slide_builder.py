"""
Construtor de apresentações PowerPoint profissionais com estrutura estratégica.
Implementa slides focados com 1 ideia por slide e design institucional.
"""

import logging
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Any
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, ChartData

from analytics.tendencias import analisar_tendencia
from config import DATA_DIR, MUNICIPIO, UF
from database import get_timeseries, list_indicators
from reports.structure import SlideType, create_presentation_structure, BRAND_COLORS
from reports.text_engine import TextGenerator, TrendAnalyzer, analyze_multiple_indicators
from reports.charts import ChartGenerator
from reports.slide_narratives import narrative_generator

logger = logging.getLogger(__name__)

TITULO_SECRETARIA = "Secretaria Municipal de Desenvolvimento, Ciência, Tecnologia e Inovação"

class SlideBuilder:
    """Construtor profissional de slides com design institucional."""
    
    def __init__(self):
        """Inicializa o construtor com configurações padrão."""
        self.prs = None
        self.chart_generator = ChartGenerator()
        self.text_generator = TextGenerator()
        self.charts_dir = DATA_DIR / "slides_charts"
        self.charts_dir.mkdir(exist_ok=True)
        
        # Configuração de cores
        self.colors = {
            "primary": RGBColor(30, 58, 138),      # Blue 900
            "secondary": RGBColor(59, 130, 246),    # Blue 500
            "success": RGBColor(22, 163, 74),      # Green 600
            "warning": RGBColor(217, 119, 6),      # Orange 600
            "error": RGBColor(220, 38, 38),        # Red 600
            "neutral": RGBColor(100, 116, 139),    # Slate 500
            "background": RGBColor(248, 250, 252),  # Slate 50
            "text": RGBColor(15, 23, 42),           # Slate 900
            "white": RGBColor(255, 255, 255)
        }
    
    def _add_logo_and_footer(self, slide):
        """Adiciona logo da prefeitura e rodapé institucional."""
        # Logo no canto superior direito
        logo_path = Path(__file__).parent.parent / "assets" / "logo_prefeitura.png"
        if logo_path.exists():
            slide.shapes.add_picture(str(logo_path), 
                                   self.prs.slide_width - Inches(1.8), 
                                   Inches(0.2), 
                                   width=Inches(1.5))
        
        # Rodapé institucional
        footer_y = self.prs.slide_height - Inches(0.5)
        tx_box = slide.shapes.add_textbox(Inches(0.5), footer_y, Inches(6), Inches(0.4))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Observatório Socioeconômico {MUNICIPIO}/{UF} | {datetime.now().year}"
        p.font.size = Pt(10)
        p.font.color.rgb = self.colors["neutral"]
    
    def _add_slide_title(self, slide, title: str, subtitle: str = ""):
        """Adiciona título formatado ao slide."""
        # Título principal
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.colors["primary"]
        p.alignment = PP_ALIGN.LEFT
        
        # Subtítulo se fornecido
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors["secondary"]
            p.alignment = PP_ALIGN.LEFT
    
    def _add_key_message(self, slide, message: str, position: str = "center"):
        """Adiciona mensagem-chave em destaque."""
        if position == "center":
            y_pos = Inches(3)
            height = Inches(2)
        else:
            y_pos = Inches(2)
            height = Inches(1.5)
        
        msg_box = slide.shapes.add_textbox(Inches(1), y_pos, Inches(8), height)
        tf = msg_box.text_frame
        p = tf.paragraphs[0]
        p.text = message
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.colors["primary"]
        p.alignment = PP_ALIGN.CENTER
    
    def _add_insights_bullets(self, slide, insights: List[str], y_start: float = 4.0):
        """Adiciona insights como bullets formatados."""
        y_pos = y_start
        
        for insight in insights:
            bullet_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.5))
            tf = bullet_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {insight}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors["text"]
            y_pos += 0.6

    def _add_source_note(self, slide, text: str):
        """Adiciona nota de fonte no rodapé do slide (padrão institucional)."""
        footer_y = self.prs.slide_height - Inches(0.85)
        box = slide.shapes.add_textbox(Inches(0.5), footer_y, Inches(9), Inches(0.3))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Fonte: {text}"
        p.font.size = Pt(10)
        p.font.color.rgb = self.colors["neutral"]
        p.alignment = PP_ALIGN.LEFT

    def _delete_slide(self, slide) -> None:
        """
        Remove um slide da apresentação.
        Necessário para garantir que não existam slides analíticos sem visualização baseada em dados reais.
        """
        try:
            slide_id = slide.slide_id
            slides = self.prs.slides
            slide_index = None
            for i, s in enumerate(slides):
                if s.slide_id == slide_id:
                    slide_index = i
                    break
            if slide_index is None:
                return

            slide_rId = slides._sldIdLst[slide_index].rId  # type: ignore[attr-defined]
            slides._sldIdLst.remove(slides._sldIdLst[slide_index])  # type: ignore[attr-defined]
            self.prs.part.drop_rel(slide_rId)
        except Exception as e:
            logger.warning("Não foi possível remover slide sem gráfico: %s", e)
    
    def _add_chart_to_slide(self, slide, chart_path: Optional[Path], 
                           x: float = 1.0, y: float = 2.0, 
                           width: float = 8.0, height: float = 4.0):
        """Adiciona gráfico ao slide."""
        if chart_path and chart_path.exists():
            try:
                # Verificar se é um arquivo de imagem válido antes de adicionar
                import os
                if os.path.getsize(chart_path) > 0:
                    # Verificar se é realmente uma imagem (não um arquivo de texto)
                    with open(chart_path, 'rb') as f:
                        header = f.read(8)
                        # Verificar assinaturas de arquivos de imagem comuns
                        image_signatures = [
                            b'\x89PNG',  # PNG
                            b'\xff\xd8\xff',  # JPEG
                            b'GIF87a',  # GIF
                            b'GIF89a',  # GIF
                            b'BM',  # BMP
                        ]
                        
                        is_image = any(header.startswith(sig) for sig in image_signatures)
                        
                        if is_image:
                            slide.shapes.add_picture(str(chart_path), Inches(x), Inches(y), 
                                                       width=Inches(width), height=Inches(height))
                        else:
                            logger.warning(f"Arquivo não é uma imagem válida: {chart_path}")
                            self._add_placeholder_text(slide, x, y, width, height)
                else:
                    logger.warning(f"Arquivo de gráfico vazio: {chart_path}")
                    self._add_placeholder_text(slide, x, y, width, height)
                    
            except Exception as e:
                logger.error(f"Erro ao adicionar gráfico ao slide: {e}")
                self._add_placeholder_text(slide, x, y, width, height)
                
        elif chart_path and hasattr(chart_path, 'getvalue'):
            # Handle BytesIO object - converter para arquivo temporário
            try:
                # Salvar BytesIO temporariamente
                temp_path = self.charts_dir / f"temp_chart_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
                with open(temp_path, 'wb') as f:
                    f.write(chart_path.getvalue())
                
                # Adicionar imagem do arquivo temporário
                slide.shapes.add_picture(str(temp_path), Inches(x), Inches(y), 
                                   width=Inches(width), height=Inches(height))
                
                # Limpar arquivo temporário após uso
                temp_path.unlink()
                
            except Exception as e:
                logger.error(f"Erro ao processar BytesIO no slide: {e}")
                self._add_placeholder_text(slide, x, y, width, height)
        else:
            # Se não houver gráfico, adicionar texto informativo
            self._add_placeholder_text(slide, x, y, width, height)
    
    def _add_placeholder_text(self, slide, x: float, y: float, width: float, height: float):
        """Adiciona nota de ausência de dados reais quando não há gráfico válido.

        Importante: não são usados dados simulados. A mensagem explicita apenas que,
        para o indicador em questão, não há série histórica suficiente ou arquivo
        de gráfico gerado a partir de dados oficiais.
        """
        tx_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
        tf = tx_box.text_frame
        p = tf.add_paragraph()
        p.text = "Sem dados oficiais suficientes para gerar o gráfico deste indicador"
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    
    def _create_cover_slide(self):
        """Cria slide de capa institucional."""
        slide_layout = self.prs.slide_layouts[6]  # blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Fundo azul institucional
        bg_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                        self.prs.slide_width, self.prs.slide_height)
        bg_rect.fill.solid()
        bg_rect.fill.fore_color.rgb = self.colors["primary"]
        bg_rect.line.fill.background()
        
        # Título principal
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), 
                                            self.prs.slide_width - Inches(1), Inches(2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = TITULO_SECRETARIA
        p.font.size = Pt(36)
        p.font.color.rgb = self.colors["white"]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Subtítulo
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), 
                                              self.prs.slide_width - Inches(1), Inches(1.5))
        tf = subtitle_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Resumo Executivo Socioeconômico\n{MUNICIPIO}/{UF}"
        p.font.size = Pt(28)
        p.font.color.rgb = self.colors["white"]
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = f"Série Analisada: 2018 a {datetime.now().year}"
        p2.font.size = Pt(20)
        p2.font.color.rgb = self.colors["secondary"]
        p2.alignment = PP_ALIGN.CENTER
        
        self._add_logo_and_footer(slide)
        return slide
    
    def _create_executive_summary_slide(self, indicators_data: Dict[str, pd.DataFrame], 
                                     ano_inicio: int, ano_fim: int):
        """Cria slide de resumo executivo com narrativa curta e ranking."""
        slide_layout = self.prs.slide_layouts[6]  # blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Título
        self._add_slide_title(slide, "Resumo Executivo", "Mensagem-chave estratégica")
        
        # Narrativa principal
        narrative = narrative_generator.generate_executive_summary(
            indicators_data, ano_inicio, ano_fim
        )
        
        # Adicionar narrativa
        tx_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.2))
        tf = tx_box.text_frame
        p = tf.add_paragraph()
        p.text = narrative
        p.font.size = Pt(18)
        p.font.color.rgb = self.colors["text"]
        p.alignment = PP_ALIGN.CENTER
        
        # Ranking automático
        best_indicators, worst_indicators = narrative_generator.generate_indicator_ranking(indicators_data)
        
        # Melhores indicadores
        if best_indicators:
            best_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(3.5), Inches(2))
            best_tf = best_box.text_frame
            best_title = best_tf.add_paragraph()
            best_title.text = "Principais avanços"
            best_title.font.size = Pt(16)
            best_title.font.bold = True
            best_title.font.color.rgb = self.colors["success"]
            
            for i, indicator in enumerate(best_indicators[:3], 1):
                p = best_tf.add_paragraph()
                p.text = f"{i}. {indicator}"
                p.font.size = Pt(14)
                p.font.color.rgb = self.colors["text"]
        
        # Piores indicadores
        if worst_indicators:
            worst_box = slide.shapes.add_textbox(Inches(5), Inches(3.2), Inches(3.5), Inches(2))
            worst_tf = worst_box.text_frame
            worst_title = worst_tf.add_paragraph()
            worst_title.text = "Pontos de atenção"
            worst_title.font.size = Pt(16)
            worst_title.font.bold = True
            worst_title.font.color.rgb = self.colors["warning"]
            
            for i, indicator in enumerate(worst_indicators[:3], 1):
                p = worst_tf.add_paragraph()
                p.text = f"{i}. {indicator}"
                p.font.size = Pt(14)
                p.font.color.rgb = self.colors["text"]
        
        # Variação percentual acumulada
        scorecard = narrative_generator.generate_scorecard_data(indicators_data)
        
        var_box = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(1))
        var_tf = var_box.text_frame
        var_title = var_tf.add_paragraph()
        var_title.text = "Variação acumulada do período"
        var_title.font.size = Pt(16)
        var_title.font.bold = True
        var_title.font.color.rgb = self.colors["primary"]
        
        # Adicionar variações principais
        var_text = ""
        for indicator, data in scorecard.items():
            if data["variation"] is not None:
                var_text += f"{indicator}: {data['variation']:+.1f}%  "
        
        if var_text:
            var_p = var_tf.add_paragraph()
            var_p.text = var_text.strip()
            var_p.font.size = Pt(14)
            var_p.font.color.rgb = self.colors["text"]
        
        self._add_logo_and_footer(slide)
        return slide
    
    def _create_key_message_slide(self, indicators_data: Dict[str, pd.DataFrame]):
        """Cria slide com mensagem-chave e 3 principais destaques."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._add_slide_title(slide, "Mensagem-Chave", "Principais destaques estratégicos")
        
        # Analisar indicadores para identificar destaques
        all_analyses = analyze_multiple_indicators(indicators_data)
        
        # Identificar 3 principais destaques
        highlights = []
        for name, analysis in all_analyses.items():
            if analysis.get("strength", 0) > 60:
                df = indicators_data.get(name)
                if df is not None and not df.empty:
                    latest = df.iloc[-1]
                    direction = analysis.get("direction", "stable")
                    
                    if direction == "increasing":
                        highlights.append(f"{name}: {latest['Valor']:,.0f} em {int(latest['Ano'])}")
                    elif direction == "decreasing":
                        highlights.append(f"{name}: {latest['Valor']:,.0f} em {int(latest['Ano'])}")
        
        # Limitar a 3 destaques mais fortes
        highlights = highlights[:3]
        
        if highlights:
            message = "3 insights essenciais para tomada de decisão:"
            self._add_key_message(slide, message)
            self._add_insights_bullets(slide, highlights, y_start=3.5)
        else:
            message = "Análise integrada dos indicadores municipais"
            self._add_key_message(slide, message)
        
        # Call-to-action
        cta_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.8))
        tf = cta_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Foco em oportunidades de crescimento sustentável"
        p.font.size = Pt(18)
        p.font.color.rgb = self.colors["success"]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        self._add_logo_and_footer(slide)
        return slide
    
    def _create_overview_slide(self, indicators_data: Dict[str, pd.DataFrame]):
        """Cria slide de panorama geral com mapa e scorecard."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Título
        self._add_slide_title(slide, "Panorama Geral", "Diagnóstico rápido do município")
        
        # Scorecard automático
        scorecard = narrative_generator.generate_scorecard_data(indicators_data)
        
        # Criar scorecard visual
        y_pos = 1.8
        for indicator_name, data in scorecard.items():
            if data["value"] is not None:
                # Caixa do indicador
                indicator_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(3), Inches(0.8))
                indicator_tf = indicator_box.text_frame
                
                # Nome do indicador
                name_p = indicator_tf.add_paragraph()
                name_p.text = indicator_name
                name_p.font.size = Pt(12)
                name_p.font.bold = True
                
                # Valor
                value_p = indicator_tf.add_paragraph()
                value_p.text = f"{data['value']:,.0f}"
                value_p.font.size = Pt(16)
                value_p.font.bold = True
                
                # Variação com cor
                var_p = indicator_tf.add_paragraph()
                if data["variation"] is not None:
                    emoji = "📈" if data["variation"] > 0 else "📉" if data["variation"] < 0 else "➡️"
                    var_color = self.colors["success"] if data["variation"] > 0 else self.colors["error"] if data["variation"] < 0 else self.colors["neutral"]
                    var_p.text = f"{emoji} {data['variation']:+.1f}%"
                    var_p.font.size = Pt(11)
                    var_p.font.color.rgb = var_color
                
                y_pos += 1.0
        
        # Mapa (informações institucionais)
        map_box = slide.shapes.add_textbox(Inches(5), Inches(1.8), Inches(4), Inches(3))
        map_tf = map_box.text_frame
        
        map_title = map_tf.add_paragraph()
        map_title.text = "Localização estratégica"
        map_title.font.size = Pt(16)
        map_title.font.bold = True
        map_title.font.color.rgb = self.colors["primary"]
        map_title.alignment = PP_ALIGN.CENTER
        
        map_desc = map_tf.add_paragraph()
        map_desc.text = "Governador Valadares\nVale do Rio Doce - MG\nPosição estratégica para logística\nAcesso a corredores econômicos"
        map_desc.font.size = Pt(14)
        map_desc.font.color.rgb = self.colors["text"]
        map_desc.alignment = PP_ALIGN.CENTER
        
        # Síntese baseada apenas em dados (quando disponível)
        msg = "Síntese: série insuficiente para síntese automatizada."
        try:
            all_analyses = analyze_multiple_indicators(indicators_data)
            pos = [k for k, a in all_analyses.items() if a.get("direction") == "increasing" and a.get("strength", 0) > 50]
            neg = [k for k, a in all_analyses.items() if a.get("direction") == "decreasing" and a.get("strength", 0) > 50]
            if pos or neg:
                parts = []
                if pos:
                    parts.append(f"crescimento em {', '.join(pos[:2])}")
                if neg:
                    parts.append(f"queda em {', '.join(neg[:2])}")
                msg = "Síntese: " + " e ".join(parts) + "."
        except Exception:
            pass

        message_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
        message_tf = message_box.text_frame
        message_p = message_tf.add_paragraph()
        message_p.text = msg
        message_p.font.size = Pt(14)
        message_p.font.bold = True
        message_p.font.color.rgb = self.colors["neutral"]
        message_p.alignment = PP_ALIGN.CENTER
        
        self._add_logo_and_footer(slide)
        return slide
                                                                        
    def _create_thematic_slide(self, theme: str, indicators_data: Dict[str, pd.DataFrame]):
        """Cria slide temático específico."""
        # Só criar slide se houver dados suficientes para gerar gráfico real
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Configurações específicas por tema
        theme_config = {
            "economia": {
                "title": "Economia",
                "subtitle": "PIB e estrutura produtiva",
                "keywords": ["pib", "vaf", "empresa"],
                "insights": ["Crescimento econômico sustentável", "Diversificação produtiva"]
            },
            "trabalho_renda": {
                "title": "Trabalho e Renda", 
                "subtitle": "Mercado formal e empreendedorismo",
                "keywords": ["emprego", "salario", "caged", "sebrae"],
                "insights": ["Geração de empregos", "Massa salarial em expansão"]
            },
            "capital_humano": {
                "title": "Capital Humano",
                "subtitle": "Educação e saúde",
                "keywords": ["matricula", "escola", "mortalidade"],
                "insights": ["Investimento em capital humano", "Qualidade de vida"]
            },
            "sustentabilidade": {
                "title": "Sustentabilidade",
                "subtitle": "Desenvolvimento sustentável",
                "keywords": ["idsc", "emissao", "area", "mapbiomas"],
                "insights": ["Equilíbrio ambiental", "Desenvolvimento consciente"]
            }
        }
        
        config = theme_config.get(theme, theme_config["economia"])
        
        self._add_slide_title(slide, config["title"], config["subtitle"])
        
        # Filtrar indicadores do tema
        theme_indicators = {}
        for name, df in indicators_data.items():
            if any(keyword in name.lower() for keyword in config["keywords"]):
                if not df.empty and len(df) > 1:
                    theme_indicators[name] = df
        
        # Adicionar gráfico principal
        chart_added = False
        if theme_indicators:
            # Selecionar o indicador mais relevante
            main_indicator = list(theme_indicators.keys())[0]
            df = theme_indicators[main_indicator]
            
            chart_path = self.charts_dir / f"{theme}_main.png"
            chart_type = "line" if len(df) > 5 else "bar"
            
            if chart_type == "line":
                self.chart_generator.create_line_chart(
                    df, "Ano", "Valor", f"Evolução de {main_indicator}",
                    output_path=chart_path
                )
            else:
                self.chart_generator.create_bar_chart(
                    df, "Ano", "Valor", f"{main_indicator} por Ano",
                    output_path=chart_path
                )
            
            self._add_chart_to_slide(slide, chart_path, y=1.5, height=3.5)
            chart_added = True
            self._add_source_note(slide, "Base de Dados Integrada (Painel GV)")
        
        # Adicionar insights
        self._add_insights_bullets(slide, config["insights"], y_start=5.5)
        
        self._add_logo_and_footer(slide)
        # Se não houve gráfico, remover o slide criado (não deixar slide analítico sem visualização)
        if not chart_added:
            self._delete_slide(slide)
            return None
        return slide
    
    def _create_trends_slide(self, indicators_data: Dict[str, pd.DataFrame]):
        """Cria slide de tendências e projeções."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._add_slide_title(slide, "Tendências e Projeções", "Cenários futuros e oportunidades")
        
        # Analisar tendências fortes
        all_analyses = analyze_multiple_indicators(indicators_data)
        strong_trends = []
        
        for name, analysis in all_analyses.items():
            if analysis.get("strength", 0) > 70:
                direction = analysis.get("direction", "stable")
                if direction in ["increasing", "decreasing"]:
                    df = indicators_data.get(name)
                    if df is not None and not df.empty:
                        strong_trends.append((name, analysis, df))
        
        # Criar gráfico com tendências mais fortes
        if strong_trends:
            # Ordenar por força
            strong_trends.sort(key=lambda x: x[1]["strength"], reverse=True)
            
            # Pegar as 3 tendências mais fortes
            top_trends = strong_trends[:3]
            trend_data = {}
            
            for name, analysis, df in top_trends:
                trend_data[name] = df
            
            chart_path = self.charts_dir / "trends_forecast.png"
            self.chart_generator.create_comparison_chart(
                trend_data, "Ano", "Valor",
                "Principais Tendências Identificadas",
                output_path=chart_path
            )
            
            self._add_chart_to_slide(slide, chart_path, y=1.5, height=3.5)
            self._add_source_note(slide, "Base de Dados Integrada (Painel GV)")
        else:
            # Sem dados suficientes -> remover slide (evita slide analítico sem visualização)
            self._delete_slide(slide)
            return None
        
        # Insights sobre tendências
        insights = [
            "Tendências de médio prazo baseadas em dados históricos",
            "Projeções estatísticas com intervalos de confiança",
            "Cenários otimista e pessimista considerados"
        ]
        
        self._add_insights_bullets(slide, insights, y_start=5.5)
        
        # Call-to-action
        cta_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.6))
        tf = cta_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Preparar o futuro hoje"
        p.font.size = Pt(18)
        p.font.color.rgb = self.colors["warning"]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        self._add_logo_and_footer(slide)
        return slide
    
    def _create_opportunities_slide(self, indicators_data: Dict[str, pd.DataFrame]):
        """Cria slide de oportunidades estratégicas."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._add_slide_title(slide, "Oportunidades Estratégicas", "Áreas prioritárias de investimento")
        
        # Analisar oportunidades
        all_analyses = analyze_multiple_indicators(indicators_data)
        
        # Identificar oportunidades (tendências positivas fortes)
        opportunities = []
        for name, analysis in all_analyses.items():
            if (analysis.get("direction") == "increasing" and 
                analysis.get("strength", 0) > 60 and
                analysis.get("confidence", 0) > 0.5):
                
                df = indicators_data.get(name)
                if df is not None and not df.empty:
                    latest = df.iloc[-1]
                    unit = str(latest.get("Unidade", "")).strip()
                    unit_txt = f" {unit}" if unit else ""
                    opportunities.append(f"{name}: {latest['Valor']:,.0f}{unit_txt} em {int(latest['Ano'])}")
        
        # Visual institucional (tabela textual) sem valores simulados
        if opportunities:
            box = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(8), Inches(3.8))
            tf = box.text_frame
            title_p = tf.paragraphs[0]
            title_p.text = "Oportunidades identificadas (com base em tendências reais)"
            title_p.font.size = Pt(16)
            title_p.font.bold = True
            for item in opportunities[:6]:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(14)
            self._add_source_note(slide, "Base de Dados Integrada (Painel GV)")
        
        # Insights sobre oportunidades
        insights = [
            "Setores com maior potencial de crescimento",
            "Investimentos prioritários para desenvolvimento",
            "Parcerias estratégicas a serem exploradas"
        ]
        
        self._add_insights_bullets(slide, insights, y_start=5.5)
        
        # Call-to-action
        cta_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.6))
        tf = cta_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Agir nas oportunidades identificadas"
        p.font.size = Pt(18)
        p.font.color.rgb = self.colors["success"]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        self._add_logo_and_footer(slide)
        return slide
    
    def _create_closing_slide(self):
        """Cria slide de encerramento."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._add_slide_title(slide, "Próximos Passos", "Agenda estratégica")
        
        # Mensagem principal
        message = "Transformar dados em ação"
        self._add_key_message(slide, message)
        
        # Próximos passos
        next_steps = [
            "Implementar plano de desenvolvimento integrado",
            "Monitorar indicadores-chave continuamente",
            "Fortalecer capacidade técnica municipal",
            "Promover transparência e participação social",
            "Avaliar impactos das intervenções"
        ]
        
        self._add_insights_bullets(slide, next_steps, y_start=3.5)
        
        # Contato/encerramento
        contact_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.8))
        tf = contact_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Observatório Socioeconômico de {MUNICIPIO}\ncontato@{MUNICIPIO}.mg.gov.br"
        p.font.size = Pt(16)
        p.font.color.rgb = self.colors["neutral"]
        p.alignment = PP_ALIGN.CENTER
        
        self._add_logo_and_footer(slide)
        return slide
    
    def build_complete_presentation(self, ano_inicio: int, ano_fim: int,
                                 output_path: Optional[str] = None) -> Path:
        """
        Constrói apresentação completa com estrutura estratégica.
        
        Args:
            ano_inicio: Ano inicial da análise
            ano_fim: Ano final da análise
            output_path: Caminho para salvar a apresentação
            
        Returns:
            Path do arquivo gerado
        """
        output_path = Path(output_path) if output_path else DATA_DIR / f"apresentacao_estrategica_{ano_inicio}_{ano_fim}.pptx"
        output_path = output_path.resolve()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        logger.info("Iniciando construção da apresentação estratégica")
        
        # Inicializar apresentação
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Coletar dados dos indicadores
        indicators_data = self._collect_indicators_data(ano_inicio, ano_fim)
        
        # Construir slides
        self._create_cover_slide()
        self._create_key_message_slide(indicators_data)
        self._create_overview_slide(indicators_data)
        self._create_thematic_slide("economia", indicators_data)
        self._create_thematic_slide("trabalho_renda", indicators_data)
        self._create_thematic_slide("capital_humano", indicators_data)
        self._create_thematic_slide("sustentabilidade", indicators_data)
        self._create_trends_slide(indicators_data)
        self._create_opportunities_slide(indicators_data)
        self._create_closing_slide()
        
        # Salvar apresentação
        self.prs.save(str(output_path))
        logger.info(f"Apresentação estratégica gerada em {output_path}")
        
        return output_path
    
    def _collect_indicators_data(self, ano_inicio: int, ano_fim: int) -> Dict[str, pd.DataFrame]:
        """Coleta dados de todos os indicadores para análise."""
        indicators_data = {}
        all_indicators = list_indicators()
        education_keys = {
            "MATRICULAS_TOTAL",
            "ESCOLAS_FUNDAMENTAL",
            "IDEB_ANOS_INICIAIS",
            "IDEB_ANOS_FINAIS",
            "TAXA_APROVACAO_FUNDAMENTAL",
        }
        
        for ind in all_indicators:
            key = ind["indicator_key"]
            source = ind["source"]

            # EDUCAÇÃO: usar exclusivamente dados provenientes de arquivos reais em data/raw
            if key in education_keys and source != "INEP_RAW":
                continue
            
            try:
                df = get_timeseries(key, source=source)
                if not df.empty:
                    # Filtrar por período
                    df_filtered = df[(df["Ano"] >= ano_inicio) & (df["Ano"] <= ano_fim)]
                    if not df_filtered.empty:
                        indicators_data[key] = df_filtered
            except Exception as e:
                logger.warning(f"Erro ao coletar indicador {key}: {e}")
                continue
        
        return indicators_data

# Função de conveniência para manter compatibilidade
def gerar_apresentacao_ppt(
    ano_inicio: int,
    ano_fim: int,
    output_path: str | Path | None = None,
) -> Path:
    """
    Gera apresentação PowerPoint estratégica completa (função de compatibilidade).
    
    Args:
        ano_inicio: Ano inicial da análise
        ano_fim: Ano final da análise
        output_path: Caminho para salvar a apresentação
        
    Returns:
        Path do arquivo gerado
    """
    builder = SlideBuilder()
    return builder.build_complete_presentation(ano_inicio, ano_fim, output_path)
